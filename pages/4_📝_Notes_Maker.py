import streamlit as st
import pdfplumber
from docx import Document
import pptx
import os
from groq import Groq

# Initialize Groq client with your API key
client = Groq(api_key="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")

# Function to load custom CSS and apply the UI theme
def load_css():
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;600;700&display=swap');
    
    .stApp {
        background-color: #1e1e1e;
        font-family: 'Poppins', sans-serif;
    }
    .main-title {
        color: #ff6347;
        font-size: 70px;
        font-weight: bold;
        text-align: center;
        margin-top: 20px;
        text-shadow: 3px 3px 6px rgba(0, 0, 0, 0.4);
    }
    .tagline {
        color: #ff4500;
        font-size: 28px;
        font-style: italic;
        text-align: center;
        margin-bottom: 30px;
    }
    .logo-container {
        text-align: center;
        margin: 20px 0;
    }
    .section-header {
        color: #ff6347;
        font-size: 36px;
        font-weight: bold;
        margin-top: 40px;
        margin-bottom: 20px;
        text-align: center;
    }
    .feature-box {
        background-color: #2f2f2f;
        border-radius: 15px;
        padding: 20px;
        margin: 20px auto;
        width: 90%;
        max-width: 700px;
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.4);
    }
    .feature-box:hover {
        background-color: #ff6347;
        transform: scale(1.05);
    }
    .feature-title {
        color: #fff;
        font-size: 22px;
        font-weight: bold;
        margin-bottom: 10px;
    }
    .feature-description {
        color: #ddd;
        font-size: 16px;
        margin-bottom: 10px;
    }
    .app-button {
        background-color: #ff6347;
        color: #fff;
        font-size: 18px;
        font-weight: bold;
        padding: 12px 24px;
        border-radius: 25px;
        margin-top: 20px;
        text-align: center;
        transition: all 0.3s ease;
        display: inline-block;
        text-decoration: none;
    }
    .app-button:hover {
        background-color: #ff4500;
        box-shadow: 0 4px 8px rgba(0,0,0,0.2);
    }
    .file-card {
        background-color: #3a3a3a;
        border-radius: 10px;
        padding: 15px;
        margin: 15px 0;
        color: #fff;
        font-weight: bold;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.3);
    }
    .file-card:hover {
        background-color: #ff6347;
    }
    .intro-statement {
        background-color: #2f2f2f;
        border-radius: 10px;
        padding: 20px;
        margin: 30px 0;
        text-align: center;
        font-size: 18px;
        color: #fff;
    }
    p {
        color: #ddd;
    }
    .progress-bar {
        margin-top: 20px;
    }
    </style>
    """, unsafe_allow_html=True)

# Function to extract text from PDFs using pdfplumber
def extract_pdf_text(uploaded_file):
    with pdfplumber.open(uploaded_file) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

# Function to extract text from DOCX files
def extract_docx_text(uploaded_file):
    doc = Document(uploaded_file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from PPTX files
def extract_pptx_text(uploaded_file):
    presentation = pptx.Presentation(uploaded_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to generate detailed notes using Groq API
def generate_detailed_notes_with_groq(text):
    chat_completion = client.chat.completions.create(
        messages=[{
            "role": "user",
            "content": f"Please summarize the following text in a structured manner with headings, subheadings, and bullet points. Make the explanation clear and understandable:\n\n{text}",
        }], 
        model="llama3-8b-8192",
    )
    return chat_completion.choices[0].message.content

# Function to extract headings and detailed content from the uploaded text
def extract_headings_and_content(text):
    headings_and_content = []
    lines = text.split('\n')
    
    current_heading = None
    current_content = []
    
    for line in lines:
        if line.strip():
            if line.isupper():
                if current_heading:
                    headings_and_content.append((current_heading, '\n'.join(current_content)))
                current_heading = line.strip()
                current_content = []
            else:
                current_content.append(line.strip())
    
    if current_heading:
        headings_and_content.append((current_heading, '\n'.join(current_content)))
    
    return headings_and_content

# Main app function
def main():
    load_css()
    
    st.title("COMSATS Chatbot 🎓")
    st.subheader("Generate Detailed Study Notes from Your Uploaded Lectures 📚")
    st.write("Upload your PDF, DOCX, or PPTX files to get detailed notes for preparation.")
    
    uploaded_files = st.file_uploader("Choose Files", type=["pdf", "docx", "pptx"], accept_multiple_files=True)
    
    if uploaded_files:
        for uploaded_file in uploaded_files:
            st.markdown(f"""
                <div class="file-card">
                    <b>{uploaded_file.name}</b> ({uploaded_file.type})
                </div>
                """, unsafe_allow_html=True)
        
        generate_button = st.button("Generate Notes 📄")
        
        if generate_button:
            full_text = ""
            
            with st.spinner('Generating notes...'):
                for uploaded_file in uploaded_files:
                    if uploaded_file.type == "application/pdf":
                        full_text += extract_pdf_text(uploaded_file)
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                        full_text += extract_docx_text(uploaded_file)
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                        full_text += extract_pptx_text(uploaded_file)
            
            if full_text:
                detailed_notes = generate_detailed_notes_with_groq(full_text)
                headings_and_content = extract_headings_and_content(full_text)
                
                st.subheader("Extracted Notes 📚")
                for heading, content in headings_and_content:
                    st.write(f"### {heading}")
                    st.write(content)
                
                st.subheader("Structured AI-Generated Notes 🤖")
                st.write(detailed_notes)
                
                st.write("For additional resources, refer to [Groq API Docs](https://groq.com/docs).")

if __name__ == "__main__":
    st.set_page_config(
        page_title="COMSATS Chatbot | Notes Maker",
        page_icon="🎓",
        layout="wide",  
        initial_sidebar_state="collapsed"
    )
    main()
