import streamlit as st
import pdfplumber
from docx import Document
import pptx
import os
from groq import Groq

# Set page config as the first command
st.set_page_config(
    page_title="Study Notes Generator | Notes Maker",
    page_icon="📚",
    layout="wide",  
    initial_sidebar_state="collapsed"
)

# Initialize Groq client with your API key
client = Groq(api_key="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")

# Function to load Poppins font
def load_css():
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;600;700&display=swap');
    
    .stApp {
        font-family: 'Poppins', sans-serif;
    }
    </style>
    """, unsafe_allow_html=True)

# Function to extract text from PDFs using pdfplumber
def extract_pdf_text(uploaded_file):
    with pdfplumber.open(uploaded_file) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

# Function to extract text from DOCX files
def extract_docx_text(uploaded_file):
    doc = Document(uploaded_file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from PPTX files
def extract_pptx_text(uploaded_file):
    presentation = pptx.Presentation(uploaded_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to generate detailed notes using Groq API
def generate_detailed_notes_with_groq(text):
    chat_completion = client.chat.completions.create(
        messages=[{
            "role": "user",
            "content": f"Please summarize the following text in a structured manner with headings, subheadings, and bullet points. Make the explanation clear and understandable:\n\n{text}",
        }], 
        model="llama3-8b-8192",
    )
    return chat_completion.choices[0].message.content

# Function to extract headings and detailed content from the uploaded text
def extract_headings_and_content(text):
    headings_and_content = []
    lines = text.split('\n')
    
    current_heading = None
    current_content = []
    
    for line in lines:
        if line.strip():
            if line.isupper():
                if current_heading:
                    headings_and_content.append((current_heading, '\n'.join(current_content)))
                current_heading = line.strip()
                current_content = []
            else:
                current_content.append(line.strip())
    
    if current_heading:
        headings_and_content.append((current_heading, '\n'.join(current_content)))
    
    return headings_and_content

# Main app function
def main():
    load_css()  # Apply Poppins font

    st.title("Study Notes Generator 📚")
    st.subheader("Generate Detailed Study Notes from Your Uploaded Lectures 📄✨")
    st.write("Upload your PDF, DOCX, or PPTX files to generate detailed and structured study notes for effective learning.")

    uploaded_files = st.file_uploader("Choose Files", type=["pdf", "docx", "pptx"], accept_multiple_files=True)

    if uploaded_files:
        for uploaded_file in uploaded_files:
            st.markdown(f"""
                <div style="padding: 10px; background-color: #f5f5f5; margin: 10px 0; border-radius: 5px;">
                    <b>{uploaded_file.name}</b> ({uploaded_file.type})
                </div>
                """, unsafe_allow_html=True)

        generate_button = st.button("Generate Notes 📄")

        if generate_button:
            full_text = ""

            with st.spinner('Generating notes...'):
                for uploaded_file in uploaded_files:
                    if uploaded_file.type == "application/pdf":
                        full_text += extract_pdf_text(uploaded_file)
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                        full_text += extract_docx_text(uploaded_file)
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                        full_text += extract_pptx_text(uploaded_file)

            if full_text:
                detailed_notes = generate_detailed_notes_with_groq(full_text)
                headings_and_content = extract_headings_and_content(full_text)
                
                st.subheader("Extracted Notes 📚")
                for heading, content in headings_and_content:
                    st.write(f"### {heading}")
                    st.write(content)
                
                st.subheader("Structured AI-Generated Notes 🤖")
                st.write(detailed_notes)
                
                st.write("For additional resources, refer to [Groq API Docs](https://groq.com/docs).")

# About the Model
st.markdown("""
    ## About This Model 🤖
    This Study Notes Generator uses AI to extract and summarize lecture materials into clear and structured study notes. It can process PDF, DOCX, and PPTX files, offering a detailed summary with headings, subheadings, and bullet points for better learning.
""")
# Footer
st.markdown(
    '<p style="text-align: center; font-weight: 600; font-size: 16px;">💻 Developed with ❤️ using Streamlit | © 2024</p>',
    unsafe_allow_html=True)

if __name__ == "__main__":
    main()
