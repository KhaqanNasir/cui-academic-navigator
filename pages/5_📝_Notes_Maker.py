import streamlit as st
import pdfplumber
from docx import Document
import pptx
import os
from groq import Groq

# Initialize Groq client with your API key
client = Groq(api_key="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")

# Function to load custom CSS and apply the UI theme
def load_css():
    st.markdown("""
    <style>
    /* Your custom CSS code here */
    </style>
    """, unsafe_allow_html=True)

# Function to extract text from PDFs using pdfplumber
def extract_pdf_text(uploaded_file):
    with pdfplumber.open(uploaded_file) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

# Function to extract text from DOCX files
def extract_docx_text(uploaded_file):
    doc = Document(uploaded_file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from PPTX files
def extract_pptx_text(uploaded_file):
    presentation = pptx.Presentation(uploaded_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to generate detailed notes using Groq API
def generate_detailed_notes_with_groq(text):
    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": f"Please summarize the following text in a structured manner with headings, subheadings, and bullet points. Make the explanation clear and understandable:\n\n{text}",
            }
        ],
        model="llama3-8b-8192",
    )
    return chat_completion.choices[0].message.content

# Function to extract headings and detailed content from the uploaded text
def extract_headings_and_content(text):
    headings_and_content = []
    lines = text.split('\n')
    
    # Simple rule to detect headings (You can adjust this based on your file content)
    current_heading = None
    current_content = []
    
    for line in lines:
        if line.strip():  # Skip empty lines
            if line.isupper():  # Assuming headings are in uppercase
                if current_heading:
                    headings_and_content.append((current_heading, '\n'.join(current_content)))
                current_heading = line.strip()
                current_content = []
            else:
                current_content.append(line.strip())
    
    # Add the last section
    if current_heading:
        headings_and_content.append((current_heading, '\n'.join(current_content)))
    
    return headings_and_content

# Main app function
def main():
    load_css()
    
    st.title("COMSATS Chatbot 🎓")
    st.subheader("Generate Detailed Study Notes from Your Uploaded Lectures 📚")
    st.write("Upload your PDF, DOCX, or PPTX files to get detailed notes for preparation.")
    
    uploaded_files = st.file_uploader("Choose Files", type=["pdf", "docx", "pptx"], accept_multiple_files=True)
    
    if uploaded_files:
        for uploaded_file in uploaded_files:
            st.write(f"Uploaded file: {uploaded_file.name}")
        
        generate_button = st.button("Generate Notes 📄")
        
        if generate_button:
            full_text = ""
            
            # Extracting text from all uploaded files
            for uploaded_file in uploaded_files:
                if uploaded_file.type == "application/pdf":
                    full_text += extract_pdf_text(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    full_text += extract_docx_text(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    full_text += extract_pptx_text(uploaded_file)
            
            if full_text:
                # Generate detailed notes using Groq API
                detailed_notes = generate_detailed_notes_with_groq(full_text)
                
                # Extract headings and content
                headings_and_content = extract_headings_and_content(full_text)
                
                # Display headings and detailed content
                st.subheader("Extracted Notes 📚")
                for heading, content in headings_and_content:
                    st.write(f"### {heading}")
                    st.write(content)
                
                # Display Groq generated notes
                st.subheader("Structured AI-Generated Notes 🤖")
                st.write(detailed_notes)
                
                # You can add links to additional resources here if needed
                st.write("For additional resources, refer to [Groq API Docs](https://groq.com/docs).")

if __name__ == "__main__":
    st.set_page_config(
        page_title="COMSATS Chatbot",
        page_icon="🎓",
        layout="wide",  
        initial_sidebar_state="collapsed"
    )
    main()
