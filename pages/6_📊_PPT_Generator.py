import openai
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import streamlit as st
from groq import Groq

# Retrieve API keys securely from Streamlit secrets
GROQ_API_KEY = st.secrets["GROQ_API_KEY4"]
OPENAI_API_KEY = st.secrets["OPENAI_API_KEY"]

# Configure OpenAI API with the key from Streamlit secrets
openai.api_key = OPENAI_API_KEY

# Initialize Groq client for Llama3 model using the API key from Streamlit secrets
client = Groq(api_key=GROQ_API_KEY)

# Streamlit configuration
st.set_page_config(
    page_title="AI-Powered Presentation Generator",
    layout="wide",  # Use the wide layout for better user experience
    page_icon="📊",
)

# Function to generate structured slide content using Llama3 model
def generate_slide_content_with_llama(topic, description, num_slides):
    """Generate slide content using the Llama3 model through the Groq API."""
    try:
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": f"Create a professional presentation with {num_slides} slides. "
                               f"Topic: {topic}. Description: {description}. Provide clear titles and bullet points for each slide."
                }
            ],
            model="llama3-8b-8192",
        )
        return response.choices[0].message.content
    except Exception as e:
        raise RuntimeError(f"Error generating slide content: {e}")


# Function to create a PowerPoint presentation with enhanced formatting
def create_presentation(slide_content, topic):
    """Create a PowerPoint presentation with provided slide content and professional formatting."""
    prs = Presentation()

    # Ensure content is well-structured
    slides = slide_content.split("Slide")
    if len(slides) <= 1:
        raise ValueError("Invalid slide content structure received from the model.")

    for slide in slides[1:]:
        slide_data = slide.strip().split("\n", 1)
        if len(slide_data) < 2:
            continue  # Skip slides without title or content

        title, content = slide_data[0].strip(), slide_data[1].strip()

        # Add title and content to the slide
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide_obj = prs.slides.add_slide(slide_layout)

        # Apply dark theme background
        background = slide_obj.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(34, 34, 34)  # Dark gray background

        # Format slide title
        slide_obj.shapes.title.text = title
        title_frame = slide_obj.shapes.title.text_frame
        title_frame.paragraphs[0].font.name = "Poppins"
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.size = Pt(30)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text

        # Add bullet points with formatting
        content_placeholder = slide_obj.placeholders[1]
        content_text_frame = content_placeholder.text_frame
        content_text_frame.clear()  # Remove default text

        for point in content.split("\n"):
            if point.strip():
                p = content_text_frame.add_paragraph()
                p.text = point.strip()
                p.font.name = "Poppins"
                p.font.size = Pt(18)  # Font size for bullet points
                p.font.color.rgb = RGBColor(200, 200, 200)  # Light gray text
                p.space_before = Pt(6)  # Add spacing between bullets

    # Save the presentation
    output_file = "generated_presentation.pptx"
    prs.save(output_file)
    return output_file


# Streamlit application
def main():
    st.markdown(
        """
        <style>
            @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700&display=swap');
            html, body, [class*="css"] {
                font-family: 'Poppins', sans-serif;
            }
        </style>
        """,
        unsafe_allow_html=True,
    )

    st.title("📊 AI-Powered Presentation Generator")
    st.markdown("Generate professional slides in minutes with the power of AI.")

    # Input options
    topic = st.text_input("🔑 Enter Topic:", placeholder="E.g., Artificial Intelligence")
    description = st.text_area("📝 Enter Description:", placeholder="Provide a detailed description of your topic.", height=200)
    num_slides = st.number_input("🔢 Number of Slides:", min_value=1, step=1)

    # Step 1: Generate content
    if st.button("Generate Content", help="Click to generate the best content for your slides"):
        if not topic or not description:
            st.error("Please enter both a topic and description.")
            return

        with st.spinner("Generating content..."):
            try:
                # Generate slide content using the Llama3 model
                slide_content = generate_slide_content_with_llama(topic, description, int(num_slides))
                st.session_state["generated_content"] = slide_content

                # Display the generated content for review
                st.subheader("📄 Generated Content")
                st.text_area("🔍 Review the Generated Slide Content:", slide_content, height=400, key="content_display")
                st.success("Content generated successfully! Review the content and proceed to generate slides.")
            except Exception as e:
                st.error(f"Error: {e}")

    # Step 2: Generate slides based on the reviewed content
    if "generated_content" in st.session_state:
        if st.button("Generate Slides", help="Click to create the slides from the generated content"):
            with st.spinner("Generating slides..."):
                try:
                    reviewed_content = st.session_state["generated_content"]
                    output_file = create_presentation(reviewed_content, topic)
                    st.success("Presentation generated successfully!")

                    # Provide download link
                    with open(output_file, "rb") as file:
                        st.download_button("📥 Download Presentation", file, file_name="presentation.pptx")
                except Exception as e:
                    st.error(f"Error: {e}")
    # Footer
    st.markdown("---")
    st.markdown(
      '<p style="text-align: center; font-weight: 600; font-size: 16px;">💻 Developed with ❤️ using Streamlit | © 2025</p>',
       unsafe_allow_html=True)


if __name__ == "__main__":
    main()


